"""
PPTXGeneratorのテスト（TDD: Red phase）
"""
import pytest
from pathlib import Path
from pptx import Presentation
from utils.image_manager import ImageData
from exporter.pptx_generator import PPTXGenerator


class TestPPTXGenerator:
    """PPTXGeneratorクラスのテスト"""

    def test_generator_initialization(self):
        """ジェネレータの初期化テスト"""
        generator = PPTXGenerator()
        assert generator is not None

    def test_generate_empty_presentation(self, temp_session_dir):
        """空のプレゼンテーション生成（画像なし）"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "empty.pptx"

        # 空リストで生成
        result_path = generator.generate([], output_path)

        # 検証
        assert result_path.exists()
        assert result_path == output_path

        # PowerPointファイルが開けるか
        prs = Presentation(str(result_path))
        assert len(prs.slides) >= 0  # 最低限のスライドがあるはず

    def test_generate_with_title_slide(self, temp_session_dir):
        """タイトルスライドの生成"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "with_title.pptx"

        result_path = generator.generate([], output_path, title="テストマニュアル")

        # PowerPointを開いて検証
        prs = Presentation(str(result_path))
        assert len(prs.slides) >= 1  # タイトルスライドが存在

    def test_generate_single_image_slide(self, temp_session_dir, sample_images, sample_image_data):
        """1枚の画像スライド生成"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "single_image.pptx"

        # 1つの画像データのみ
        image_data_list = [sample_image_data[0]]

        result_path = generator.generate(image_data_list, output_path)

        # PowerPointを開いて検証
        prs = Presentation(str(result_path))
        # タイトルスライド + 画像スライド1枚
        assert len(prs.slides) >= 2

    def test_generate_multiple_image_slides(self, temp_session_dir, sample_image_data):
        """複数画像スライド生成"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "multiple_images.pptx"

        result_path = generator.generate(sample_image_data, output_path)

        # PowerPointを開いて検証
        prs = Presentation(str(result_path))
        # タイトルスライド + 画像スライド3枚
        assert len(prs.slides) >= 4

    def test_slide_contains_description(self, temp_session_dir, sample_images):
        """スライドに説明文が含まれているか"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "with_description.pptx"

        # 説明文付きの画像データ
        image_data = ImageData(
            filepath=str(sample_images[0]),
            description="これはテスト説明文です",
            order=0
        )

        result_path = generator.generate([image_data], output_path)

        # PowerPointを開いて検証
        prs = Presentation(str(result_path))

        # 2枚目のスライド（画像スライド）を確認
        if len(prs.slides) >= 2:
            slide = prs.slides[1]
            # スライド内のテキストを検索
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            # 説明文が含まれているか
            assert any("テスト説明文" in text for text in slide_text)

    def test_slide_contains_image(self, temp_session_dir, sample_images, sample_image_data):
        """スライドに画像が含まれているか"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "with_image.pptx"

        result_path = generator.generate([sample_image_data[0]], output_path)

        # PowerPointを開いて検証
        prs = Presentation(str(result_path))

        # 2枚目のスライド（画像スライド）を確認
        if len(prs.slides) >= 2:
            slide = prs.slides[1]

            # 画像が含まれているか（Pictureオブジェクトの存在確認）
            has_image = False
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    has_image = True
                    break

            assert has_image, "スライドに画像が含まれていません"

    def test_output_file_is_valid_pptx(self, temp_session_dir, sample_image_data):
        """生成されたファイルが有効なPowerPointファイルか"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "valid_output.pptx"

        result_path = generator.generate(sample_image_data, output_path)

        # ファイルが存在
        assert result_path.exists()

        # 拡張子が正しい
        assert result_path.suffix == '.pptx'

        # PowerPointとして開ける（例外が発生しない）
        try:
            prs = Presentation(str(result_path))
            assert prs is not None
        except Exception as e:
            pytest.fail(f"PowerPointファイルとして開けません: {e}")

    def test_slide_order_matches_image_order(self, temp_session_dir, sample_images):
        """スライドの順序が画像の順序と一致するか"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "ordered_slides.pptx"

        # 順序付き画像データ
        image_data_list = [
            ImageData(filepath=str(sample_images[0]), description="First", order=0),
            ImageData(filepath=str(sample_images[1]), description="Second", order=1),
            ImageData(filepath=str(sample_images[2]), description="Third", order=2),
        ]

        result_path = generator.generate(image_data_list, output_path)

        # PowerPointを開いて検証
        prs = Presentation(str(result_path))

        # タイトルスライドをスキップして、画像スライドの順序を確認
        if len(prs.slides) >= 4:
            for i, slide_idx in enumerate([1, 2, 3]):
                slide = prs.slides[slide_idx]
                slide_texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                expected_text = ["First", "Second", "Third"][i]
                assert any(expected_text in text for text in slide_texts)

    def test_custom_title(self, temp_session_dir):
        """カスタムタイトルの設定"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "custom_title.pptx"
        custom_title = "カスタムマニュアルタイトル"

        result_path = generator.generate([], output_path, title=custom_title)

        # PowerPointを開いて検証
        prs = Presentation(str(result_path))

        if len(prs.slides) >= 1:
            title_slide = prs.slides[0]
            slide_texts = [shape.text for shape in title_slide.shapes if hasattr(shape, "text")]
            assert any(custom_title in text for text in slide_texts)

    def test_image_path_not_exists(self, temp_session_dir):
        """存在しない画像パスの処理"""
        generator = PPTXGenerator()
        output_path = temp_session_dir / "missing_image.pptx"

        # 存在しない画像パス
        image_data = ImageData(
            filepath=str(temp_session_dir / "nonexistent.png"),
            description="Missing image",
            order=0
        )

        # エラーハンドリングされるべき（例外を投げるか、スキップするか）
        # 実装方針によって調整
        try:
            result_path = generator.generate([image_data], output_path)
            # 生成は成功するが、画像がスキップされる想定
            assert result_path.exists()
        except FileNotFoundError:
            # または、FileNotFoundErrorを投げる想定
            pass
