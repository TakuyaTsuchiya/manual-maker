"""
統合テスト
収録→編集→出力の完全フローをテスト
"""
import pytest
from pathlib import Path
from PIL import Image
from utils.image_manager import ImageManager
from exporter.pptx_generator import PPTXGenerator


class TestIntegration:
    """統合テストクラス"""

    def test_basic_setup(self, temp_session_dir):
        """基本的なセットアップテスト"""
        assert temp_session_dir.exists()
        assert temp_session_dir.is_dir()

    def test_session_load_and_save(self, temp_session_dir, sample_images):
        """セッションの読み込みと保存テスト"""
        # 初回ロード - 画像を自動検出
        manager1 = ImageManager(temp_session_dir)
        assert len(manager1.images) == 3

        # 説明文を追加
        manager1.update_description(0, "Test description 1")
        manager1.update_description(1, "Test description 2")
        manager1.save_metadata()

        # 新しいインスタンスでロード - 保存した説明文が読み込まれるか
        manager2 = ImageManager(temp_session_dir)
        assert len(manager2.images) == 3
        assert manager2.images[0].description == "Test description 1"
        assert manager2.images[1].description == "Test description 2"
        assert manager2.images[2].description == ""  # Empty string, not None

    def test_session_metadata_persistence(self, temp_session_dir, sample_images):
        """メタデータ永続化の詳細テスト"""
        manager = ImageManager(temp_session_dir)

        # 複数の操作を実行
        manager.update_description(0, "First")
        manager.update_description(1, "Second")
        manager.update_description(2, "Third")
        manager.save_metadata()

        # metadata.jsonが存在するか
        metadata_file = temp_session_dir / "metadata.json"
        assert metadata_file.exists()

        # 新しいマネージャーで読み込み
        new_manager = ImageManager(temp_session_dir)
        assert new_manager.images[0].description == "First"
        assert new_manager.images[1].description == "Second"
        assert new_manager.images[2].description == "Third"

        # 順序も保持されているか
        for i, img in enumerate(new_manager.images):
            assert img.order == i

    def test_multiple_undo_redo(self, temp_session_dir, sample_images):
        """複数回のUndo/Redoテスト"""
        manager = ImageManager(temp_session_dir)
        initial_count = len(manager.images)

        # 操作1: 説明文更新
        manager.update_description(0, "Step 1")
        assert manager.images[0].description == "Step 1"
        assert len(manager.undo_stack) == 1

        # 操作2: 説明文更新
        manager.update_description(1, "Step 2")
        assert manager.images[1].description == "Step 2"
        assert len(manager.undo_stack) == 2

        # 操作3: 削除
        manager.delete_image(2)
        assert len(manager.images) == initial_count - 1
        assert len(manager.undo_stack) == 3

        # Undo 1回目: 削除を取り消し
        assert manager.undo() is True
        assert len(manager.images) == initial_count
        assert len(manager.undo_stack) == 2

        # Undo 2回目: Step 2を取り消し
        assert manager.undo() is True
        assert manager.images[1].description == ""
        assert len(manager.undo_stack) == 1

        # Undo 3回目: Step 1を取り消し
        assert manager.undo() is True
        assert manager.images[0].description == ""
        assert len(manager.undo_stack) == 0

        # もうUndoできない
        assert manager.undo() is False

    def test_undo_preserves_order(self, temp_session_dir, sample_images):
        """Undoが順序を正しく復元するかテスト"""
        manager = ImageManager(temp_session_dir)

        # 元の順序を記録
        original_paths = [img.filepath for img in manager.images]

        # 並び替え
        manager.swap_images(0, 2)
        reordered_paths = [img.filepath for img in manager.images]

        # 順序が変わっているか
        assert reordered_paths != original_paths
        assert reordered_paths == [original_paths[2], original_paths[1], original_paths[0]]

        # Undo
        manager.undo()
        restored_paths = [img.filepath for img in manager.images]

        # 元の順序に戻っているか
        assert restored_paths == original_paths

    def test_full_workflow_record_edit_export(self, temp_session_dir):
        """完全フロー: 収録→編集→出力テスト"""
        # ステップ1: 収録（画像ファイル作成をシミュレート）
        from PIL import Image
        image_paths = []
        for i in range(5):
            img = Image.new('RGB', (800, 600), color=(i*50, 100, 200))
            img_path = temp_session_dir / f"screenshot_{i:04d}.png"
            img.save(img_path)
            image_paths.append(img_path)

        # ステップ2: 編集（ImageManagerで画像管理）
        manager = ImageManager(temp_session_dir)
        assert len(manager.images) == 5

        # 説明文を追加
        manager.update_description(0, "アプリケーションを起動")
        manager.update_description(1, "ファイルを開く")
        manager.update_description(2, "編集モードに切り替え")
        manager.update_description(3, "変更を保存")
        manager.update_description(4, "アプリケーションを終了")

        # 2枚目を削除
        manager.delete_image(1)
        assert len(manager.images) == 4

        # Undoして復元
        manager.undo()
        assert len(manager.images) == 5

        # 並び替え（最後と最初を入れ替え）
        manager.swap_images(0, 4)

        # ステップ3: PowerPoint出力
        generator = PPTXGenerator()
        output_path = temp_session_dir / "manual.pptx"

        result = generator.generate(
            manager.get_images(),
            output_path,
            title="統合テストマニュアル"
        )

        # 検証
        assert result.exists()
        assert result.suffix == '.pptx'

        # PowerPointを開いて確認
        from pptx import Presentation
        prs = Presentation(str(result))

        # タイトルスライド + 画像スライド5枚 = 6枚
        assert len(prs.slides) >= 6

        # タイトルスライドの内容確認
        title_slide = prs.slides[0]
        title_texts = [shape.text for shape in title_slide.shapes if hasattr(shape, "text")]
        assert any("統合テストマニュアル" in text for text in title_texts)
