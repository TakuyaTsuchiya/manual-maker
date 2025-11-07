"""
PowerPoint生成モジュール
"""
from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from utils.image_manager import ImageData


# レイアウト定数
SLIDE_LAYOUT_TITLE = 0  # タイトルスライド
SLIDE_LAYOUT_BLANK = 6  # 空白スライド

# 画像配置定数
IMAGE_LEFT = Inches(1)
IMAGE_TOP = Inches(1)
IMAGE_HEIGHT = Inches(4.5)

# 説明文配置定数
DESC_LEFT = Inches(1)
DESC_TOP = Inches(5.8)
DESC_WIDTH = Inches(8)
DESC_HEIGHT = Inches(1)
DESC_FONT_SIZE = Pt(14)

# デフォルト値
DEFAULT_TITLE = "マニュアル"


class PPTXGenerator:
    """PowerPoint生成クラス"""

    def __init__(self):
        """初期化"""
        pass

    def generate(
        self,
        image_data_list: List[ImageData],
        output_path: Path,
        title: Optional[str] = None
    ) -> Path:
        """
        PowerPointファイルを生成

        Args:
            image_data_list: 画像データのリスト
            output_path: 出力ファイルパス
            title: プレゼンテーションのタイトル（オプション）

        Returns:
            Path: 生成されたファイルのパス
        """
        # 新しいプレゼンテーションを作成
        prs = Presentation()

        # タイトルスライドを作成（タイトル指定時または画像がある場合）
        if title or len(image_data_list) > 0:
            self._create_title_slide(prs, title or DEFAULT_TITLE)

        # 画像スライドを作成
        for img_data in image_data_list:
            if Path(img_data.filepath).exists():
                self._create_content_slide(prs, img_data)

        # ファイルを保存
        prs.save(str(output_path))

        return output_path

    def _create_title_slide(self, prs: Presentation, title: str) -> None:
        """
        タイトルスライドを作成

        Args:
            prs: プレゼンテーションオブジェクト
            title: タイトル文字列
        """
        title_slide_layout = prs.slide_layouts[SLIDE_LAYOUT_TITLE]
        slide = prs.slides.add_slide(title_slide_layout)

        # タイトルを設定
        title_shape = slide.shapes.title
        title_shape.text = title

    def _create_content_slide(self, prs: Presentation, img_data: ImageData) -> None:
        """
        コンテンツスライドを作成（画像 + 説明文）

        Args:
            prs: プレゼンテーションオブジェクト
            img_data: 画像データ
        """
        blank_slide_layout = prs.slide_layouts[SLIDE_LAYOUT_BLANK]
        slide = prs.slides.add_slide(blank_slide_layout)

        # 画像を追加
        self._add_image_to_slide(slide, img_data.filepath)

        # 説明文を追加
        if img_data.description:
            self._add_description_to_slide(slide, img_data.description)

    def _add_image_to_slide(self, slide, image_path: str) -> None:
        """
        スライドに画像を追加

        Args:
            slide: スライドオブジェクト
            image_path: 画像ファイルパス
        """
        slide.shapes.add_picture(
            image_path,
            IMAGE_LEFT,
            IMAGE_TOP,
            height=IMAGE_HEIGHT
        )

    def _add_description_to_slide(self, slide, description: str) -> None:
        """
        スライドに説明文を追加

        Args:
            slide: スライドオブジェクト
            description: 説明文
        """
        textbox = slide.shapes.add_textbox(
            DESC_LEFT,
            DESC_TOP,
            DESC_WIDTH,
            DESC_HEIGHT
        )
        text_frame = textbox.text_frame
        text_frame.text = description

        # フォント設定
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = DESC_FONT_SIZE
